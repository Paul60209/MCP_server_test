# 首先修復 python-pptx 在 Python 3.10+ 中的相容性問題
import sys
import collections
if not hasattr(collections, 'Container'):
    import collections.abc
    collections.Container = collections.abc.Container
    collections.Mapping = collections.abc.Mapping
    collections.MutableMapping = collections.abc.MutableMapping
    collections.Sequence = collections.abc.Sequence
    collections.Set = collections.abc.Set

# 現在引入所需模組
import asyncio
import tempfile
import os
import time
import dotenv
import base64
import pathlib
import argparse
import json

# 引入 OpenAI API
from langchain_openai import ChatOpenAI

# 引入 MCP 相關模組
from mcp.server.fastmcp import FastMCP, Image

# 最後再引入 python-pptx 相關模組
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_THEME_COLOR_INDEX, MSO_COLOR_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Pt

# 載入環境變數
dotenv.load_dotenv()

# 定義輸出路徑
OUTPUT_PATH = 'output'

# 創建 MCP 伺服器
mcp = FastMCP("PPTTranslatorServer")

# 解析命令行參數
parser = argparse.ArgumentParser(description='PPT翻譯MCP服務器')
parser.add_argument('--port', type=int, default=8003, help='服務器監聽端口 (默認: 8003)')
args = parser.parse_args()

# 設置環境變數，讓 FastMCP 使用指定的端口
os.environ["MCP_SSE_PORT"] = str(args.port)

def get_text_frame_properties(text_frame):
    """獲取文本框的所有格式屬性"""
    properties = {
        'margin_left': text_frame.margin_left,
        'margin_right': text_frame.margin_right,
        'margin_top': text_frame.margin_top,
        'margin_bottom': text_frame.margin_bottom,
        'vertical_anchor': text_frame.vertical_anchor,
        'word_wrap': text_frame.word_wrap,
        'auto_size': text_frame.auto_size,
    }
    return properties

def get_paragraph_properties(paragraph):
    """獲取段落的所有格式屬性"""
    properties = {
        'alignment': paragraph.alignment,
        'level': paragraph.level,
        'line_spacing': paragraph.line_spacing,
        'space_before': paragraph.space_before,
        'space_after': paragraph.space_after,
    }
    return properties

def get_color_properties(color):
    """獲取顏色屬性"""
    if not color:
        return None
        
    properties = {
        'type': color.type if hasattr(color, 'type') else None,
        'rgb': color.rgb if hasattr(color, 'rgb') else None,
        'theme_color': color.theme_color if hasattr(color, 'theme_color') else None,
        'brightness': color.brightness if hasattr(color, 'brightness') else None,
    }
    return properties

def get_run_properties(run):
    """獲取文本運行的所有格式屬性"""
    font = run.font
    properties = {
        'size': font.size,
        'name': font.name,
        'bold': font.bold,
        'italic': font.italic,
        'underline': font.underline,
        'color': get_color_properties(font.color),
        'fill': get_color_properties(font.fill.fore_color) if hasattr(font, 'fill') else None,
    }
    return properties

def apply_color_properties(color_obj, properties):
    """應用顏色屬性"""
    if not properties or not color_obj:
        return
        
    try:
        # 如果有 RGB 值，直接設置 RGB 顏色
        if properties['rgb']:
            if isinstance(properties['rgb'], (tuple, list)) and len(properties['rgb']) == 3:
                color_obj.rgb = RGBColor(*properties['rgb'])
            else:
                color_obj.rgb = properties['rgb']
        # 如果有主題顏色，設置主題顏色
        elif properties['theme_color'] and properties['theme_color'] != MSO_THEME_COLOR_INDEX.NOT_THEME_COLOR:
            color_obj.theme_color = properties['theme_color']
            if properties['brightness'] is not None:
                color_obj.brightness = properties['brightness']
    except Exception as e:
        print(f"設置顏色時發生錯誤: {str(e)}")
        pass  # 如果設置失敗，保持原有顏色

def apply_text_frame_properties(text_frame, properties):
    """應用文本框格式屬性"""
    text_frame.margin_left = properties['margin_left']
    text_frame.margin_right = properties['margin_right']
    text_frame.margin_top = properties['margin_top']
    text_frame.margin_bottom = properties['margin_bottom']
    text_frame.vertical_anchor = properties['vertical_anchor']
    text_frame.word_wrap = properties['word_wrap']
    text_frame.auto_size = properties['auto_size']

def apply_paragraph_properties(paragraph, properties):
    """應用段落格式屬性"""
    paragraph.alignment = properties['alignment']
    paragraph.level = properties['level']
    paragraph.line_spacing = properties['line_spacing']
    paragraph.space_before = properties['space_before']
    paragraph.space_after = properties['space_after']

def apply_run_properties(run, properties):
    """應用文本運行格式屬性"""
    font = run.font
    if properties['size']:
        font.size = properties['size']
    if properties['name']:
        font.name = properties['name']
    if properties['bold'] is not None:
        font.bold = properties['bold']
    if properties['italic'] is not None:
        font.italic = properties['italic']
    if properties['underline'] is not None:
        font.underline = properties['underline']
    
    # 應用顏色
    if properties['color']:
        apply_color_properties(font.color, properties['color'])
    if properties['fill'] and hasattr(font, 'fill'):
        apply_color_properties(font.fill.fore_color, properties['fill'])

async def translate_text(text: str, olang: str, tlang: str, ctx=None) -> str:
    """使用 ChatGPT 翻譯文本。

    Args:
        text (str): 要翻譯的文本
        olang (str): 原始語言代碼
        tlang (str): 目標語言代碼
        ctx: MCP 上下文對象 (不再使用)

    Returns:
        str: 翻譯後的文本
    """
    if not text.strip():
        return text

    print(f"\n正在翻譯文本:")
    print(f"原文 ({olang}): {text}")

    try:
        # 創建 ChatGPT 模型
        model = ChatOpenAI(temperature=0)
        
        # 創建系統提示
        system_message = f"""You are a professional translator. Translate the following text from {olang} to {tlang}.
        Rules:
        1. Keep all formatting symbols (like bullet points, numbers) unchanged
        2. Keep all special characters unchanged
        3. Keep all whitespace and line breaks
        4. Only translate the actual text content
        5. Maintain the same tone and style
        6. Do not add any explanations or notes
        7. Keep all numbers and dates unchanged
        8. Keep all proper nouns unchanged unless they have standard translations
        """
        
        # 創建消息列表
        messages = [
            {"role": "system", "content": system_message},
            {"role": "user", "content": text}
        ]
        
        # 執行翻譯
        response = await model.ainvoke(messages)
        translated_text = response.content.strip()
        
        print(f"譯文 ({tlang}): {translated_text}\n")
        return translated_text
        
    except Exception as e:
        print(f"翻譯失敗: {str(e)}")
        if ctx:
            await ctx.info(f"翻譯失敗: {str(e)}")
        # 返回原文以確保不會丟失內容
        return text

async def translate_group_shape(shape, olang: str, tlang: str, ctx=None) -> None:
    """翻譯群組中的所有形狀。

    Args:
        shape: PowerPoint 群組形狀對象
        olang (str): 原始語言代碼
        tlang (str): 目標語言代碼
        ctx: MCP 上下文對象
    """
    try:
        if not hasattr(shape, 'shapes'):
            return
            
        # 遍歷群組中的所有形狀
        for child_shape in shape.shapes:
            if child_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # 遞歸處理嵌套群組
                await translate_group_shape(child_shape, olang, tlang, ctx)
            else:
                # 翻譯單個形狀
                await translate_shape(child_shape, olang, tlang, ctx)
    except Exception as e:
        print(f"翻譯群組形狀時發生錯誤: {str(e)}")
        if ctx:
            await ctx.info(f"翻譯群組形狀時發生錯誤: {str(e)}")
        raise

async def translate_shape(shape, olang: str, tlang: str, ctx=None) -> None:
    """翻譯 PowerPoint 中的形狀。

    Args:
        shape: PowerPoint 形狀對象
        olang (str): 原始語言代碼
        tlang (str): 目標語言代碼
        ctx: MCP 上下文對象
    """
    try:
        # 處理群組形狀
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            print(f"[信息] 檢測到群組形狀，進行遞歸處理")
            await translate_group_shape(shape, olang, tlang, ctx)
            return
            
        # 檢查形狀是否包含文本框
        if not hasattr(shape, "text_frame"):
            print(f"[信息] 跳過不含文本框的形狀")
            return
            
        text_frame = shape.text_frame
        if not text_frame.text.strip():
            print(f"[信息] 跳過空文本框")
            return
            
        print(f"\n===== 處理形狀文本 =====")
        print(f"[信息] 形狀文本預覽: {text_frame.text[:50] + '...' if len(text_frame.text) > 50 else text_frame.text}")
        
        # 保存文本框格式
        text_frame_props = get_text_frame_properties(text_frame)
        
        # 遍歷所有段落
        paragraph_count = len(text_frame.paragraphs)
        print(f"[信息] 段落數量: {paragraph_count}")
        
        for i, paragraph in enumerate(text_frame.paragraphs, 1):
            # 保存段落格式
            para_props = get_paragraph_properties(paragraph)
            
            print(f"[信息] 處理第 {i}/{paragraph_count} 段落")
            
            # 遍歷所有文本運行
            runs_data = []
            run_count = len(paragraph.runs)
            
            for j, run in enumerate(paragraph.runs, 1):
                # 保存運行格式和文本
                run_props = get_run_properties(run)
                original_text = run.text
                print(f"[信息] 處理第 {j}/{run_count} 運行文本: {original_text[:30] + '...' if len(original_text) > 30 else original_text}")
                
                if original_text.strip():
                    # 使用 ChatGPT 進行翻譯
                    translated_text = await translate_text(original_text, olang, tlang, ctx)
                    runs_data.append((translated_text, run_props))
                else:
                    # 跳過空運行文本
                    runs_data.append((original_text, run_props))
            
            # 清除原有內容
            for _ in range(len(paragraph.runs)):
                paragraph._p.remove(paragraph.runs[0]._r)
            
            # 添加翻譯後的文本並應用格式
            for text, props in runs_data:
                run = paragraph.add_run()
                run.text = text
                apply_run_properties(run, props)
            
            # 恢復段落格式
            apply_paragraph_properties(paragraph, para_props)
        
        # 恢復文本框格式
        apply_text_frame_properties(text_frame, text_frame_props)
        print(f"===== 形狀文本處理完成 =====\n")
        
    except Exception as e:
        print(f"[錯誤] 翻譯形狀時發生錯誤: {str(e)}")
        if ctx:
            await ctx.info(f"翻譯形狀時發生錯誤: {str(e)}")
        raise

async def translate_ppt_file(file_path: str, olang: str, tlang: str, ctx=None) -> str:
    """翻譯 PowerPoint 文件。

    Args:
        file_path (str): PowerPoint 文件路徑
        olang (str): 原始語言代碼
        tlang (str): 目標語言代碼
        ctx: MCP 上下文對象

    Returns:
        str: 翻譯後的文件路徑
    """
    try:
        print("\n========== PowerPoint 翻譯過程 ==========")
        # 1. 建立輸出目錄
        os.makedirs(OUTPUT_PATH, exist_ok=True)
        print(f"[信息] 輸出目錄: {os.path.abspath(OUTPUT_PATH)}")
        
        # 2. 準備輸出文件路徑
        file_name = os.path.basename(file_path)
        name, ext = os.path.splitext(file_name)
        output_file = f'translated_{name}{ext}'
        output_path = os.path.join(OUTPUT_PATH, output_file)
        print(f"[信息] 輸入文件: {file_path}")
        print(f"[信息] 輸出文件: {output_path}")
        
        # 3. 載入 PowerPoint
        print(f"[信息] 開始 PowerPoint 翻譯...")
        print(f"[信息] 源語言: {olang}")
        print(f"[信息] 目標語言: {tlang}")
        if ctx:
            await ctx.info(f"開始翻譯...\n從 {olang} 到 {tlang}")
        
        print(f"[信息] 載入 PowerPoint 文件...")
        presentation = Presentation(file_path)
        total_slides = len(presentation.slides)
        print(f"[信息] 檢測到 {total_slides} 張投影片")
        
        # 4. 翻譯每個投影片
        for index, slide in enumerate(presentation.slides, 1):
            progress_msg = f"正在翻譯投影片 {index}/{total_slides}..."
            print(f"\n[進度] {progress_msg}")
            
            # 計數投影片中的形狀數量
            shape_count = len(slide.shapes)
            print(f"[信息] 投影片 {index} 包含 {shape_count} 個形狀")
            
            if ctx:
                await ctx.info(progress_msg)
                # 報告進度 (0-100%)
                await ctx.report_progress(index - 1, total_slides)
            
            # 翻譯投影片中的每個形狀
            for shape_idx, shape in enumerate(slide.shapes, 1):
                print(f"[信息] 處理投影片 {index}, 形狀 {shape_idx}/{shape_count}")
                await translate_shape(shape, olang, tlang, ctx)
        
        # 5. 儲存翻譯後的文件
        print("\n[信息] 儲存翻譯後的文件...")
        if ctx:
            await ctx.info("翻譯完成，生成文件中...")
            await ctx.report_progress(total_slides, total_slides)  # 100% 完成
            
        presentation.save(output_path)
        print(f"[成功] 翻譯後的文件已保存至: {output_path}")
        print(f"========== PowerPoint 翻譯完成 ==========\n")
        
        # 6. 返回路徑
        return output_path
        
    except Exception as e:
        error_msg = f"翻譯過程中發生錯誤: {str(e)}"
        print(f"\n[錯誤] {error_msg}")
        if ctx:
            await ctx.info(error_msg)
        raise

@mcp.tool()
async def translate_ppt(olang: str, tlang: str, file_content: str = None, file_name: str = None) -> str:
    """
    翻譯 PowerPoint 檔案從一種語言到另一種語言，同時保留原始格式。
    
    ## 使用場景
    - 需要將演示文稿翻譯成其他語言時
    - 準備多語言簡報
    - 國際會議/演講準備
    
    ## 參數說明
    :param olang: 來源語言代碼或語言名稱，例如 'zh-TW'、'繁體中文'、'english'等
    :param tlang: 目標語言代碼或語言名稱，例如 'en'、'英文'、'japanese'等
    :param file_content: PowerPoint 檔案的內容（base64 編碼字符串）
    :param file_name: 檔案名稱（可選，用於確定檔案類型）
    
    ## 輸入範例
    - 從中文到英文：olang="中文"，tlang="英文"
    - 從英文到日文：olang="english"，tlang="japanese"
    - 從日文到中文：olang="ja"，tlang="zh-TW"
    
    ## 檔案要求
    - 支援.ppt和.pptx格式
    - 檔案大小不超過10MB
    - 保留原始格式包括字體、顏色、排版等
    
    ## 注意事項
    - 翻譯大型文件可能需要幾分鐘時間
    - 複雜的圖表和特殊格式可能無法完美保留
    - 專有名詞可能需要手動修正
    
    :return: 包含翻譯結果訊息和檔案內容的 JSON 字符串
    """
    # 獲取 MCP 上下文
    try:
        ctx = mcp.get_current_request_context()
        print(f"[信息] 成功獲取 MCP 上下文")
    except Exception as e:
        print(f"[警告] 無法獲取 MCP 上下文: {str(e)}")
        ctx = None
    
    try:
        print(f"\n========== PPT翻譯工具啟動 ==========")
        print(f"[參數] 來源語言: {olang}")
        print(f"[參數] 目標語言: {tlang}")
        print(f"[參數] 檔案名稱: {file_name if file_name else '未指定'}")
        print(f"[參數] 文件內容長度: {len(file_content) if file_content else 0} 字符")
        
        # 檢查必要參數
        if not file_content:
            print(f"[錯誤] 未提供文件內容")
            return json.dumps({
                "success": False,
                "message": "錯誤：需要提供 PowerPoint 檔案內容。請上傳檔案並提供內容。"
            })
        
        # 確保檔案名稱有效
        if not file_name:
            file_name = "uploaded_presentation.pptx"
            print(f"[信息] 未提供檔案名稱，使用預設名稱: {file_name}")
        elif not (file_name.lower().endswith('.ppt') or file_name.lower().endswith('.pptx')):
            old_name = file_name
            file_name += ".pptx"  # 添加預設副檔名
            print(f"[信息] 檔案名稱 '{old_name}' 沒有有效副檔名，修改為: {file_name}")
        else:
            print(f"[信息] 使用提供的檔案名稱: {file_name}")
        
        # 建立臨時檔案來存儲上傳的內容
        print(f"[信息] 創建臨時文件...")
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file_name)[1]) as temp_file:
            # 直接寫入檔案內容
            try:
                if isinstance(file_content, bytes):
                    print(f"[信息] 檢測到二進制文件內容")
                    temp_file.write(file_content)
                elif isinstance(file_content, str):
                    print(f"[信息] 檢測到字符串文件內容，嘗試Base64解碼")
                    # 嘗試解碼 base64 字串
                    try:
                        decoded_content = base64.b64decode(file_content)
                        temp_file.write(decoded_content)
                        print(f"[信息] Base64解碼成功，寫入 {len(decoded_content)} 字節")
                    except Exception as e:
                        print(f"[警告] Base64解碼失敗: {str(e)}，嘗試寫入純文本")
                        # 如果不是有效的base64，視為純文本
                        encoded_content = file_content.encode('utf-8')
                        temp_file.write(encoded_content)
                        print(f"[信息] 寫入 {len(encoded_content)} 字節的純文本")
                else:
                    print(f"[錯誤] 不支援的文件內容類型: {type(file_content)}")
                    return json.dumps({
                        "success": False,
                        "message": "錯誤：不支援的檔案內容格式。請提供二進制或base64編碼的檔案內容。"
                    })
                
                temp_file_path = temp_file.name
                print(f"[信息] 檔案已暫存於: {temp_file_path}")
            except Exception as e:
                print(f"[錯誤] 處理檔案內容時發生錯誤: {str(e)}")
                return json.dumps({
                    "success": False,
                    "message": f"處理檔案內容時發生錯誤: {str(e)}。"
                })
        
        # 執行翻譯 - 直接使用用戶提供的語言參數，不進行驗證或轉換
        print("[信息] 開始翻譯流程...")
        output_path = await translate_ppt_file(temp_file_path, olang, tlang, ctx)
        print(f"[信息] 翻譯完成，結果路徑: {output_path}")
        
        # 清理臨時檔案
        if os.path.exists(temp_file_path):
            os.remove(temp_file_path)
            print(f"[信息] 已清理臨時檔案: {temp_file_path}")
        
        # 讀取翻譯後的文件並編碼為base64
        with open(output_path, "rb") as f:
            file_bytes = f.read()
            print(f"[信息] 讀取翻譯後的文件: {len(file_bytes)} 字節")
            translated_file_content = base64.b64encode(file_bytes).decode('utf-8')
            print(f"[信息] 編碼為Base64: {len(translated_file_content)} 字符")
        
        # 獲取輸出文件名
        output_file_name = os.path.basename(output_path)
        
        # 返回包含必要訊息的 JSON
        print(f"[信息] 準備返回結果 JSON")
        result_json = json.dumps({
            "success": True,
            "message": "翻譯完成！",
            "file_name": output_file_name,
            "file_content": translated_file_content
        })
        print(f"[信息] 結果 JSON 長度: {len(result_json)} 字符")
        print(f"========== PPT翻譯工具完成 ==========\n")
        return result_json
        
    except Exception as e:
        print(f"[錯誤] 翻譯工具執行錯誤: {str(e)}")
        error_result = json.dumps({
            "success": False,
            "message": f"翻譯過程中發生錯誤: {str(e)}"
        })
        print(f"========== PPT翻譯工具出錯 ==========\n")
        return error_result

@mcp.resource("translator://instructions")
async def get_instructions() -> str:
    """獲取 PPT 翻譯器的使用說明。"""
    return """
# PowerPoint 翻譯工具使用說明

這個工具可以將 PowerPoint 文件從一種語言翻譯成另一種語言，同時保留原始格式。

## 支援的語言

- 中文（zh-TW）
- 英文（en）
- 日文（ja）
- 其他語言代碼也可嘗試

## 使用方法

1. 將 PowerPoint 檔案（.ppt 或 .pptx）轉為 base64 編碼
2. 調用 `translate_ppt` 工具，提供以下參數：
   - olang: 原始語言代碼
   - tlang: 目標語言代碼
   - file_content: 檔案的 base64 編碼內容

## 翻譯過程

1. 工具會解析 PowerPoint 檔案中的所有文本
2. 逐一翻譯每個文本元素，保留原始格式
3. 生成新的 PowerPoint 檔案，將結果回傳

## 注意事項

- 翻譯大型檔案可能需要較長時間
- 某些複雜的格式可能無法完美保留
- 檔案大小限制為 10MB
"""

# 啟動伺服器
if __name__ == "__main__":
    # 確保輸出目錄存在
    os.makedirs(OUTPUT_PATH, exist_ok=True)
    
    import uvicorn
    from starlette.applications import Starlette
    from starlette.routing import Route, Mount
    from mcp.server.sse import SseServerTransport
    
    # 啟動 MCP 伺服器
    print(f"啟動PPT翻譯服務器 在端口 {args.port}")
    
    # 創建 SSE 傳輸
    sse = SseServerTransport("/mcp/")
    
    # 定義 SSE 連接處理函數
    async def handle_sse(request):
        async with sse.connect_sse(
            request.scope, request.receive, request._send
        ) as streams:
            await mcp._mcp_server.run(
                streams[0], 
                streams[1],
                mcp._mcp_server.create_initialization_options()
            )
    
    # 定義健康檢查端點
    async def health_check(request):
        from starlette.responses import JSONResponse
        return JSONResponse({"status": "ok"})
    
    # 創建 Starlette 應用
    starlette_app = Starlette(
        routes=[
            # SSE 端點
            Route("/sse", endpoint=handle_sse, methods=["GET"]),
            Mount("/mcp/", app=sse.handle_post_message),
            # 添加健康檢查端點
            Route("/health", endpoint=health_check, methods=["GET"]),
        ]
    )
    
    # 使用 uvicorn 啟動應用
    uvicorn.run(starlette_app, host="0.0.0.0", port=args.port) 